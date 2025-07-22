"""
Document loader for RAG-Ultra: efficiently extracts per-page text from PDFs, DOCX, PPTX, and TXT files.
Supports lazy loading, parallel processing, and robust fallback strategies for different file types.
"""

import os
from typing import Dict, Union, Optional, Any, Iterator, Tuple
from concurrent.futures import ThreadPoolExecutor, as_completed

from .config import SUPPORTED_FILE_TYPES
from .utils import logger

# Optional dependency imports at module level
try:
    import fitz  # PyMuPDF
    FITZ_AVAILABLE = True
except ImportError:
    FITZ_AVAILABLE = False

try:
    import PyPDF2
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False

try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import pptx
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Default number of workers for parallel processing
DEFAULT_MAX_WORKERS = 8

def convert_document_to_text(document_path: str, max_workers: int = DEFAULT_MAX_WORKERS) -> Dict[int, str]:
    """
    Extract per-page text from a document, auto-detecting file type.

    Args:
        document_path: Path to the document file.
        max_workers: Maximum number of worker threads for parallel processing.
    Returns:
        Dict mapping page numbers (1-indexed) to extracted text.
    Raises:
        ValueError: If document type is not supported.
        FileNotFoundError: If document file doesn't exist.
    """
    if not os.path.exists(document_path):
        raise FileNotFoundError(f"Document not found at {document_path}")
    file_ext = os.path.splitext(document_path)[1].lower().lstrip('.')
    if file_ext not in SUPPORTED_FILE_TYPES:
        raise ValueError(f"Unsupported file type: {file_ext}. Supported types: {SUPPORTED_FILE_TYPES}")
    logger.info(f"Processing {file_ext} document: {document_path}")
    if file_ext == "pdf":
        return extract_text_from_pdf(document_path, max_workers)
    elif file_ext == "docx":
        return extract_text_from_docx(document_path)
    elif file_ext == "pptx":
        return extract_text_from_pptx(document_path)
    elif file_ext == "txt":
        return extract_text_from_txt(document_path)
    else:
        raise ValueError(f"Unsupported file type: {file_ext}")

def lazy_extract_text(document_path: str) -> Iterator[Tuple[int, str]]:
    """
    Lazily extract text from a document, yielding one page at a time.
    Useful for very large documents to avoid loading everything into memory.

    Args:
        document_path: Path to the document file.
    Yields:
        Tuples of (page_number, page_text).
    """
    if not os.path.exists(document_path):
        raise FileNotFoundError(f"Document not found at {document_path}")
    file_ext = os.path.splitext(document_path)[1].lower().lstrip('.')
    if file_ext not in SUPPORTED_FILE_TYPES:
        raise ValueError(f"Unsupported file type: {file_ext}. Supported types: {SUPPORTED_FILE_TYPES}")
    logger.info(f"Lazy processing {file_ext} document: {document_path}")
    if file_ext == "pdf":
        if FITZ_AVAILABLE:
            doc = fitz.open(document_path)
            for page_idx in range(len(doc)):
                try:
                    page = doc[page_idx]
                    yield (page_idx + 1, page.get_text())
                except Exception as e:
                    logger.error(f"Error extracting text from PDF page {page_idx+1}: {e}")
                    yield (page_idx + 1, "")
            doc.close()
        elif PYPDF2_AVAILABLE:
            with open(document_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                for i in range(len(reader.pages)):
                    try:
                        page = reader.pages[i]
                        yield (i + 1, page.extract_text())
                    except Exception as e:
                        logger.error(f"Error extracting text from PDF page {i+1}: {e}")
                        yield (i + 1, "")
        else:
            raise ImportError("No PDF extraction library available. Install PyMuPDF or PyPDF2.")
    elif file_ext == "docx":
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx library is required for DOCX processing")
        doc = docx.Document(document_path)
        current_page = 1
        current_text = ""
        for para in doc.paragraphs:
            if "PAGE BREAK" in para.text.upper() or para.text.strip() == "\f":
                if current_text.strip():
                    yield (current_page, current_text.strip())
                    current_page += 1
                    current_text = ""
            else:
                current_text += para.text + "\n"
        if current_text.strip():
            yield (current_page, current_text.strip())
    elif file_ext == "pptx":
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx library is required for PPTX processing")
        presentation = pptx.Presentation(document_path)
        for i, slide in enumerate(presentation.slides):
            text_list = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_list.append(shape.text)
            yield (i + 1, "\n".join(text_list))
    elif file_ext == "txt":
        with open(document_path, 'r', encoding='utf-8') as file:
            content = file.read()
        pages = split_text_into_pages(content)
        for i, page in enumerate(pages, 1):
            if page.strip():
                yield (i, page.strip())

def _extract_pdf_page_fitz(pdf_path: str, page_idx: int) -> Tuple[int, str]:
    """
    Helper for parallel PDF extraction using PyMuPDF.
    Returns (1-indexed page number, text).
    """
    try:
        doc = fitz.open(pdf_path)
        page = doc[page_idx]
        text = page.get_text()
        doc.close()
        return (page_idx + 1, text)
    except Exception as e:
        logger.error(f"Error extracting text from PDF page {page_idx+1}: {e}")
        return (page_idx + 1, "")

def _extract_pdf_page_pypdf2(pdf_path: str, page_idx: int) -> Tuple[int, str]:
    """
    Helper for parallel PDF extraction using PyPDF2.
    Returns (1-indexed page number, text).
    """
    try:
        with open(pdf_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            page = reader.pages[page_idx]
            return (page_idx + 1, page.extract_text())
    except Exception as e:
        logger.error(f"Error extracting text from PDF page {page_idx+1}: {e}")
        return (page_idx + 1, "")

def extract_text_from_pdf(pdf_path: str, max_workers: int = DEFAULT_MAX_WORKERS) -> Dict[int, str]:
    """
    Extract text from each page of a PDF using parallel processing.
    Prefers PyMuPDF for speed, falls back to PyPDF2 if needed.

    Args:
        pdf_path: Path to the PDF file.
        max_workers: Maximum number of worker threads.
    Returns:
        Dict mapping page numbers to page text.
    Raises:
        ImportError: If no PDF extraction library is available.
    """
    if FITZ_AVAILABLE:
        try:
            doc = fitz.open(pdf_path)
            page_count = len(doc)
            doc.close()
            page_text = {}
            actual_workers = min(max_workers, max(1, page_count))
            with ThreadPoolExecutor(max_workers=actual_workers) as executor:
                future_to_page = {
                    executor.submit(_extract_pdf_page_fitz, pdf_path, page_idx): page_idx
                    for page_idx in range(page_count)
                }
                for future in as_completed(future_to_page):
                    page_num, text = future.result()
                    page_text[page_num] = text
            logger.info(f"Extracted text from {len(page_text)} PDF pages using PyMuPDF")
            return page_text
        except Exception as e:
            logger.warning(f"Error using PyMuPDF for PDF text extraction: {e}. Falling back to PyPDF2.")
    if PYPDF2_AVAILABLE:
        try:
            with open(pdf_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                page_count = len(reader.pages)
            page_text = {}
            actual_workers = min(max_workers, max(1, page_count))
            with ThreadPoolExecutor(max_workers=actual_workers) as executor:
                future_to_page = {
                    executor.submit(_extract_pdf_page_pypdf2, pdf_path, page_idx): page_idx
                    for page_idx in range(page_count)
                }
                for future in as_completed(future_to_page):
                    page_num, text = future.result()
                    page_text[page_num] = text
            logger.info(f"Extracted text from {len(page_text)} PDF pages using PyPDF2")
            return page_text
        except Exception as e:
            logger.error(f"Error extracting text from PDF: {e}")
            raise
    logger.error("No PDF library available. Install PyMuPDF or PyPDF2.")
    raise ImportError("No PDF library available. Install PyMuPDF or PyPDF2.")

def extract_text_from_docx(docx_path: str) -> Dict[int, str]:
    """
    Extract text from a DOCX file, splitting by page breaks.
    Each section between page breaks is treated as a "page".

    Args:
        docx_path: Path to the Word document.
    Returns:
        Dict mapping simulated page numbers to text.
    Raises:
        ImportError: If python-docx is not installed.
    """
    if not DOCX_AVAILABLE:
        logger.error("python-docx library is required for DOCX processing. Install with 'pip install python-docx'")
        raise ImportError("python-docx library is required for DOCX processing")
    page_text = {}
    try:
        doc = docx.Document(docx_path)
        current_page = 1
        current_text = ""
        for para in doc.paragraphs:
            if ("PAGE BREAK" in para.text.upper() or
                para.text.strip() == "\f" or
                (hasattr(para, "style") and para.style and "page break" in str(para.style).lower())):
                if current_text.strip():
                    page_text[current_page] = current_text.strip()
                    current_page += 1
                    current_text = ""
            else:
                current_text += para.text + "\n"
        if current_text.strip():
            page_text[current_page] = current_text.strip()
        logger.info(f"Extracted {len(page_text)} pages from DOCX document")
    except Exception as e:
        logger.error(f"Error extracting text from DOCX: {e}")
        raise
    return page_text

def extract_text_from_pptx(pptx_path: str) -> Dict[int, str]:
    """
    Extract text from a PPTX file, treating each slide as a page.

    Args:
        pptx_path: Path to the PowerPoint presentation.
    Returns:
        Dict mapping slide numbers to slide text.
    Raises:
        ImportError: If python-pptx is not installed.
    """
    if not PPTX_AVAILABLE:
        logger.error("python-pptx library is required for PPTX processing. Install with 'pip install python-pptx'")
        raise ImportError("python-pptx library is required for PPTX processing")
    page_text = {}
    try:
        presentation = pptx.Presentation(pptx_path)
        for i, slide in enumerate(presentation.slides):
            text_list = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_list.append(shape.text)
            page_text[i + 1] = "\n".join(text_list)
        logger.info(f"Extracted text from {len(page_text)} slides")
    except Exception as e:
        logger.error(f"Error extracting text from PPTX: {e}")
        raise
    return page_text

def split_text_into_pages(content: str, max_chars_per_page: int = 3000) -> list:
    """
    Split text into pages using common page markers or by paragraph length.
    Tries to preserve natural breaks and keep page sizes reasonable.

    Args:
        content: Text content to split.
        max_chars_per_page: Maximum characters per page.
    Returns:
        List of page content strings.
    """
    page_markers = ["\f", "----", "****", "======", "# Page", "===", "---", "***"]
    for marker in page_markers:
        if marker in content:
            pages = content.split(marker)
            logger.info(f"Split text file by marker: {marker}")
            return pages
    paragraphs = content.split("\n\n")
    pages = []
    current_page = ""
    for para in paragraphs:
        if len(current_page) + len(para) > max_chars_per_page and current_page:
            pages.append(current_page)
            current_page = para
        else:
            if current_page:
                current_page += "\n\n" + para
            else:
                current_page = para
    if current_page:
        pages.append(current_page)
    logger.info(f"Split text file into {len(pages)} pages by paragraph breaks")
    return pages

def extract_text_from_txt(txt_path: str) -> Dict[int, str]:
    """
    Extract text from a TXT file, splitting by page markers or length.

    Args:
        txt_path: Path to the text file.
    Returns:
        Dict mapping page numbers to page text.
    Raises:
        Exception: If file cannot be read or split.
    """
    try:
        with open(txt_path, 'r', encoding='utf-8') as file:
            content = file.read()
        pages = split_text_into_pages(content)
        page_text = {i+1: page.strip() for i, page in enumerate(pages) if page.strip()}
        return page_text
    except Exception as e:
        logger.error(f"Error extracting text from TXT: {e}")
        raise
